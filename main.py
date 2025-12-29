import json
import os
import tempfile
from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.dml import MSO_THEME_COLOR
from werkzeug.utils import secure_filename

app = Flask(__name__)

# Configuration
app.config['MAX_CONTENT_LENGTH'] = int(os.getenv('MAX_CONTENT_LENGTH', 16 * 1024 * 1024))  # 16MB default
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'dev-secret-key-change-in-production')
FONT = "Calibri"

# Helper function to parse color
def parse_color(color_spec):
    """Parse color from various formats: 'red', '#FF0000', or [255, 0, 0]"""
    if isinstance(color_spec, list) and len(color_spec) == 3:
        return RGBColor(*color_spec)
    elif isinstance(color_spec, str):
        if color_spec.startswith('#'):
            # Hex color
            color_spec = color_spec.lstrip('#')
            return RGBColor(*[int(color_spec[i:i+2], 16) for i in (0, 2, 4)])
        else:
            # Named colors
            color_map = {
                'red': RGBColor(255, 0, 0),
                'blue': RGBColor(0, 0, 255),
                'green': RGBColor(0, 128, 0),
                'yellow': RGBColor(255, 255, 0),
                'orange': RGBColor(255, 165, 0),
                'purple': RGBColor(128, 0, 128),
                'pink': RGBColor(255, 192, 203),
                'brown': RGBColor(165, 42, 42),
                'gray': RGBColor(128, 128, 128),
                'black': RGBColor(0, 0, 0),
                'white': RGBColor(255, 255, 255),
            }
            return color_map.get(color_spec.lower(), RGBColor(0, 0, 0))
    return RGBColor(0, 0, 0)  # Default black

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/health')
def health():
    """Health check endpoint for monitoring"""
    return jsonify({
        'status': 'healthy',
        'service': 'ppt-generator',
        'version': '1.0.0'
    })



def style(run, size=18, bold=False, italic=False, color=None, underline=False):
    """Enhanced style function with more formatting options"""
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.underline = underline
    if color:
        run.font.color.rgb = parse_color(color)
    else:
        run.font.color.rgb = RGBColor(0, 0, 0)

def add_text(tf, text, level=0, size=18, bold=False, italic=False, color=None, align=None):
    """Enhanced text function with alignment and formatting"""
    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
    p.text = text
    p.level = level
    if align:
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        p.alignment = alignment_map.get(align.lower(), PP_ALIGN.LEFT)
    for r in p.runs:
        style(r, size, bold, italic, color)

def render_heading(tf, block):
    """Render heading block"""
    text = block.get("text", "")
    level = block.get("level", 1)  # 1=large, 2=medium, 3=small
    color = block.get("color")
    align = block.get("align", "left")
    
    size_map = {1: 28, 2: 24, 3: 20}
    size = size_map.get(level, 24)
    
    add_text(tf, text, level=0, size=size, bold=True, color=color, align=align)

def render_paragraph(tf, block):
    """Render styled paragraph"""
    text = block.get("text", "")
    size = block.get("size", 18)
    color = block.get("color")
    bold = block.get("bold", False)
    italic = block.get("italic", False)
    align = block.get("align", "left")
    
    add_text(tf, text, size=size, bold=bold, italic=italic, color=color, align=align)

def render_bullets(tf, block):
    """Enhanced bullet rendering with styles"""
    items = block.get("items", [])
    bullet_color = block.get("color")
    
    for item in items:
        if isinstance(item, str):
            add_text(tf, item, color=bullet_color)
        else:
            # Styled bullet
            text = item.get("text", "")
            bold = item.get("bold", True)
            color = item.get("color", bullet_color)
            add_text(tf, text, bold=bold, color=color)
            
            # Sub-points
            for sub in item.get("subpoints", []):
                if isinstance(sub, str):
                    add_text(tf, sub, level=1, color=bullet_color)
                else:
                    sub_text = sub.get("text", "")
                    sub_color = sub.get("color", bullet_color)
                    add_text(tf, sub_text, level=1, color=sub_color)

def render_table(slide, block, top=2.5):
    """Render a table on the slide"""
    rows = block.get("rows", [])
    if not rows:
        return
    
    num_rows = len(rows)
    num_cols = len(rows[0]) if rows else 0
    
    # Position and size
    left = Inches(block.get("left", 1))
    top = Inches(block.get("top", top))
    width = Inches(block.get("width", 8))
    height = Inches(block.get("height", 0.4 * num_rows))
    
    # Create table
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table
    
    # Set column widths
    col_widths = block.get("col_widths", [width.inches / num_cols] * num_cols)
    for i, col_width in enumerate(col_widths[:num_cols]):
        table.columns[i].width = Inches(col_width)
    
    # Fill table
    header = block.get("header", True)
    header_color = parse_color(block.get("header_color", [68, 114, 196]))  # Blue
    header_text_color = parse_color(block.get("header_text_color", "white"))
    
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            
            # Cell text
            if isinstance(cell_data, dict):
                cell.text = str(cell_data.get("text", ""))
                cell_color = cell_data.get("color")
                bg_color = cell_data.get("bg_color")
            else:
                cell.text = str(cell_data)
                cell_color = None
                bg_color = None
            
            # Format cell
            cell.text_frame.paragraphs[0].font.size = Pt(block.get("font_size", 11))
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Header row styling
            if header and row_idx == 0:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = header_text_color
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif bg_color:
                cell.fill.solid()
                cell.fill.fore_color.rgb = parse_color(bg_color)
            
            if cell_color:
                cell.text_frame.paragraphs[0].font.color.rgb = parse_color(cell_color)

def render_text_box(slide, block):
    """Render a text box with custom positioning and styling"""
    text = block.get("text", "")
    left = Inches(block.get("left", 1))
    top = Inches(block.get("top", 3))
    width = Inches(block.get("width", 6))
    height = Inches(block.get("height", 1))
    
    # Create text box
    text_box = slide.shapes.add_textbox(left, top, width, height)
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    
    # Styling
    font_size = block.get("font_size", 18)
    bold = block.get("bold", False)
    italic = block.get("italic", False)
    color = block.get("color")
    bg_color = block.get("bg_color")
    align = block.get("align", "left")
    
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    if color:
        p.font.color.rgb = parse_color(color)
    
    # Alignment
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    p.alignment = alignment_map.get(align.lower(), PP_ALIGN.LEFT)
    
    # Background color
    if bg_color:
        text_box.fill.solid()
        text_box.fill.fore_color.rgb = parse_color(bg_color)

def render_numbered_list(tf, block):
    """Render numbered list"""
    items = block.get("items", [])
    start = block.get("start", 1)
    color = block.get("color")
    
    for i, item in enumerate(items, start=start):
        if isinstance(item, str):
            add_text(tf, f"{i}. {item}", color=color)
        else:
            text = item.get("text", "")
            item_color = item.get("color", color)
            add_text(tf, f"{i}. {text}", color=item_color)

def render_images(slide, block, top=2.5):
    items = block["items"]
    count = len(items)

    if block["layout"] == "row":
        width = 8 / count
        for i, img in enumerate(items):
            try:
                slide.shapes.add_picture(
                    img["path"],
                    Inches(0.5 + i * width),
                    Inches(top),
                    width=Inches(width - 0.2)
                )
            except FileNotFoundError:
                print(f"Warning: Image file not found at {img['path']}")

    elif block["layout"] == "column":
        height = 4 / count
        for i, img in enumerate(items):
            try:
                slide.shapes.add_picture(
                    img["path"],
                    Inches(2),
                    Inches(top + i * height),
                    height=Inches(height - 0.2)
                )
            except FileNotFoundError:
                print(f"Warning: Image file not found at {img['path']}")

    elif block["layout"] == "grid":
        cols = 2
        rows = (count + 1) // 2
        w, h = 4, 2
        for idx, img in enumerate(items):
            try:
                r, c = divmod(idx, cols)
                slide.shapes.add_picture(
                    img["path"],
                    Inches(0.5 + c * w),
                    Inches(top + r * h),
                    width=Inches(w - 0.3)
                )
            except FileNotFoundError:
                print(f"Warning: Image file not found at {img['path']}")

def render_blocks(slide, blocks):
    """Enhanced block renderer supporting multiple content types"""
    # Check if we have a text placeholder
    has_text_placeholder = len(slide.placeholders) > 1
    
    if has_text_placeholder:
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()
    else:
        tf = None
    
    current_top = 2.5  # Starting position for non-text elements
    
    for block in blocks:
        kind = block.get("kind", "")
        
        # Text-based blocks (use placeholder text frame)
        if kind == "heading" and tf:
            render_heading(tf, block)
            
        elif kind == "paragraph" and tf:
            render_paragraph(tf, block)
            
        elif kind == "bullets" and tf:
            render_bullets(tf, block)
            
        elif kind == "numbered_list" and tf:
            render_numbered_list(tf, block)
        
        # Positioned blocks (custom positioning)
        elif kind == "table":
            render_table(slide, block, current_top)
            current_top += block.get("height", 2) + 0.3
            
        elif kind == "text_box":
            render_text_box(slide, block)
            
        elif kind == "images":
            render_images(slide, block, current_top)
            # Calculate height used by images
            if block.get("layout") == "column":
                current_top += block.get("height", 4)
            else:
                current_top += 2.5
        
        # Legacy support for old "paragraph" without explicit kind
        elif "text" in block and kind == "" and tf:
            add_text(tf, block["text"])


def create_college_title_slide(prs, college_data):
    """Create custom college/university title slide"""
    # Use blank layout for custom design
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add college logo at top center (if exists)
    logo_path = 'static/college.png'
    if os.path.exists(logo_path):
        left = Inches(3.5)  # Center position
        top = Inches(0.5)
        height = Inches(1.2)
        slide.shapes.add_picture(logo_path, left, top, height=height)
    
    # Add college/university name
    uni_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.5))
    uni_tf = uni_box.text_frame
    uni_tf.text = college_data.get('college_name', 'College/University Name')
    uni_p = uni_tf.paragraphs[0]
    uni_p.alignment = PP_ALIGN.CENTER
    uni_p.font.size = Pt(24)
    uni_p.font.bold = True
    uni_p.font.color.rgb = RGBColor(0, 0, 128)
    
    # Add presentation title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(8), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = college_data.get('title', 'Presentation Title')
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Check if single or group
    if college_data.get('type') == 'single':
        # Single student details
        details_top = 3.8
        details_box = slide.shapes.add_textbox(Inches(2), Inches(details_top), Inches(6), Inches(2.5))
        details_tf = details_box.text_frame
        details_tf.word_wrap = True
        
        # Submitted by
        p = details_tf.paragraphs[0]
        p.text = "Submitted By:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Student details
        student_info = [
            f"Name: {college_data.get('student_name', '')}",
            f"USN: {college_data.get('usn', '')}",
            f"Course: {college_data.get('course', '')}",
            f"Semester: {college_data.get('semester', '')}"
        ]
        
        for info in student_info:
            p = details_tf.add_paragraph()
            p.text = info
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(6)
        
        # Submitted to
        p = details_tf.add_paragraph()
        p.text = f"\nSubmitted To: {college_data.get('professor', '')}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(12)
        
    elif college_data.get('type') == 'group':
        # Group - create table for students
        students = college_data.get('students', [])
        
        # Add "Submitted By:" text
        sub_box = slide.shapes.add_textbox(Inches(2), Inches(3.7), Inches(6), Inches(0.4))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = "Submitted By:"
        sub_p.font.size = Pt(16)
        sub_p.font.bold = True
        sub_p.alignment = PP_ALIGN.CENTER
        
        # Create table
        if students:
            rows = len(students) + 1  # +1 for header
            cols = 2
            left = Inches(2.5)
            top = Inches(4.2)
            width = Inches(5)
            height = Inches(0.4 * rows)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set column widths
            table.columns[0].width = Inches(2.5)
            table.columns[1].width = Inches(2.5)
            
            # Header row
            table.cell(0, 0).text = "Name"
            table.cell(0, 1).text = "USN"
            
            for i, cell in enumerate([table.cell(0, 0), table.cell(0, 1)]):
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
            
            # Data rows
            for idx, student in enumerate(students):
                table.cell(idx + 1, 0).text = student.get('name', '')
                table.cell(idx + 1, 1).text = student.get('usn', '')
                
                for col in range(2):
                    cell = table.cell(idx + 1, col)
                    cell.text_frame.paragraphs[0].font.size = Pt(11)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Submitted to
        prof_box = slide.shapes.add_textbox(Inches(2), Inches(6.5), Inches(6), Inches(0.4))
        prof_tf = prof_box.text_frame
        prof_p = prof_tf.paragraphs[0]
        prof_p.text = f"Submitted To: {college_data.get('professor', '')}"
        prof_p.font.size = Pt(14)
        prof_p.font.bold = True
        prof_p.alignment = PP_ALIGN.CENTER


@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    try:
        # Get JSON data from request or use default content.json
        if request.is_json:
            data = request.get_json()
            if 'json_data' in data:
                # Parse JSON string from frontend
                content = json.loads(data['json_data'])
                file_name = data.get('file_name', 'presentation')
                jain_data = data.get('jain_data')  # Get college/university data if provided
            else:
                content = data
                file_name = 'presentation'
                jain_data = None
        else:
            # Fallback to content.json file
            with open('content.json', 'r', encoding='utf-8') as f:
                content = json.load(f)
            file_name = 'Generated'
            jain_data = None
        
        # Validate required fields
        if 'meta' not in content or 'slides' not in content:
            return jsonify({'error': 'Invalid JSON structure. Required: meta and slides'}), 400
        
        # Create presentation
        prs = Presentation()

        # Add college title slide if requested
        if jain_data and jain_data.get('enabled'):
            create_college_title_slide(prs, jain_data)
        else:
            # Standard title slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = content["meta"].get("title", "Presentation")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content["meta"].get("subtitle", "")

        # Content slides
        for s in content["slides"]:
            if s.get("type") == "title":
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = s.get("title", "Slide")

            if "subtitle" in s:
                p = slide.shapes.title.text_frame.add_paragraph()
                p.text = s["subtitle"]
                p.level = 1

            render_blocks(slide, s.get("blocks", []))

            if "notes" in s:
                slide.notes_slide.notes_text_frame.text = s["notes"]

        # Save to temp directory (writable on serverless platforms like Vercel)
        temp_dir = tempfile.gettempdir()
        output_filename = f"{secure_filename(file_name)}.pptx"
        output_path = os.path.join(temp_dir, output_filename)
        
        prs.save(output_path)
        print(f"✅ PPT generated: {output_path}")
        
        # Send file and let Flask clean up after response
        return send_file(output_path, 
                        as_attachment=True, 
                        download_name=output_filename,
                        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except json.JSONDecodeError as e:
        return jsonify({'error': f'Invalid JSON: {str(e)}'}), 400
    except KeyError as e:
        return jsonify({'error': f'Missing required field: {str(e)}'}), 400
    except Exception as e:
        print(f"Error generating PPT: {str(e)}")
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    # Development server only
    port = int(os.getenv('PORT', 5000))
    debug = os.getenv('FLASK_ENV', 'development') == 'development'
    app.run(debug=debug, host='0.0.0.0', port=port)
